import os, sys, io, json, re, time, math, requests
from datetime import datetime, timedelta
from concurrent.futures import ThreadPoolExecutor, as_completed
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from urllib.parse import quote
from staticmap import StaticMap, CircleMarker, Line

IMG_DIR = "img"
MAP_DIR = "maps"
DATA_FILE = "trip_data.json"
OUTPUT_FILE = "Iceland_Trip.pptx"
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

os.makedirs(IMG_DIR, exist_ok=True)
os.makedirs(MAP_DIR, exist_ok=True)

USER_AGENT = "IcelandTripPlanner/2.0"

# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def sanitize_filename(name):
    s = name.replace(" ", "_").replace("/", "_")
    s = re.sub(r"[^\w.()\-]", "", s)
    if not s:
        s = "image"
    return s


def _download_image(url, save_path):
    ir = requests.get(url, headers={"User-Agent": USER_AGENT}, timeout=15)
    if ir.status_code == 429:
        return 429
    if ir.status_code != 200:
        return ir.status_code
    ct = ir.headers.get("Content-Type", "")
    if not ct.startswith("image"):
        return None
    ext = ".jpg"
    if "jpeg" in ct or "jpg" in ct:
        ext = ".jpg"
    elif "png" in ct:
        ext = ".png"
    elif "webp" in ct:
        ext = ".webp"
    else:
        return None
    path = save_path + ext
    with open(path, "wb") as f:
        f.write(ir.content)
    return path


def _ascii_fold(name):
    replacements = {
        'á': 'a', 'ð': 'd', 'é': 'e', 'í': 'i', 'ó': 'o', 'ú': 'u',
        'ý': 'y', 'þ': 'th', 'æ': 'ae', 'ö': 'o',
        'Á': 'A', 'Ð': 'D', 'É': 'E', 'Í': 'I', 'Ó': 'O', 'Ú': 'U',
        'Ý': 'Y', 'Þ': 'Th', 'Æ': 'Ae', 'Ö': 'O',
    }
    for accented, ascii_rep in replacements.items():
        name = name.replace(accented, ascii_rep)
    return name


def _search_wikipedia_title(query):
    for attempt in range(2):
        try:
            url = "https://en.wikipedia.org/w/api.php"
            params = {
                "action": "opensearch",
                "search": query,
                "limit": 3,
                "namespace": 0,
                "format": "json",
            }
            r = requests.get(url, params=params, headers={"User-Agent": USER_AGENT}, timeout=15)
            if r.status_code == 429:
                time.sleep(3 * (attempt + 1))
                continue
            r.raise_for_status()
            data = r.json()
            titles = data[1] if len(data) > 1 else []
            if titles:
                return titles[0]
            return None
        except Exception:
            if attempt < 2:
                time.sleep(2)
                continue
            return None
    return None


def _fetch_from_wikipedia(wiki_title, location_name):
    titles_to_try = [wiki_title]

    searched = _search_wikipedia_title(wiki_title)
    if searched and searched.lower() != wiki_title.lower():
        titles_to_try.append(searched)

    ascii_name = _ascii_fold(wiki_title)
    if ascii_name.lower() != wiki_title.lower():
        searched_ascii = _search_wikipedia_title(ascii_name)
        if searched_ascii and searched_ascii.lower() not in (t.lower() for t in titles_to_try):
            titles_to_try.append(searched_ascii)

    for attempt_title in titles_to_try:
        for attempt in range(2):
            try:
                url = "https://en.wikipedia.org/api/rest_v1/page/summary/" + quote(attempt_title)
                r = requests.get(url, headers={"User-Agent": USER_AGENT}, timeout=15)
                if r.status_code == 429:
                    time.sleep(3 * (attempt + 1))
                    continue
                if r.status_code == 404:
                    break
                r.raise_for_status()
                data = r.json()
                thumb = data.get("thumbnail", {}).get("source")
                if not thumb:
                    break
                base = sanitize_filename(attempt_title)
                save_path = os.path.join(IMG_DIR, base)
                result = _download_image(thumb, save_path)
                if result == 429:
                    time.sleep(3 * (attempt + 1))
                    continue
                if isinstance(result, str):
                    return result, "wikipedia"
                if result is None:
                    break
                return None, f"http {result}"
            except requests.exceptions.HTTPError as e:
                if e.response is not None and e.response.status_code == 404:
                    break
                if attempt < 2:
                    time.sleep(2 * (attempt + 1))
                    continue
                return None, str(e)
            except Exception as e:
                if attempt < 2:
                    time.sleep(2 * (attempt + 1))
                    continue
                return None, f"{type(e).__name__}: {e}"
    return None, "page not found"


def _fetch_from_commons(wiki_title):
    queries = [wiki_title, _ascii_fold(wiki_title), wiki_title + " Iceland"]
    for attempt in range(min(len(queries), 3)):
        query = queries[attempt]
        try:
            search_url = "https://commons.wikimedia.org/w/api.php"
            params = {
                "action": "query",
                "list": "search",
                "srsearch": query,
                "srnamespace": 6,
                "format": "json",
                "srlimit": 5,
            }
            r = requests.get(search_url, params=params, headers={"User-Agent": USER_AGENT}, timeout=15)
            if r.status_code == 429:
                time.sleep(3)
                continue
            r.raise_for_status()
            data = r.json()

            file_titles = [res["title"] for res in data.get("query", {}).get("search", [])
                           if res["title"].startswith("File:")]
            if not file_titles:
                continue

            params = {
                "action": "query",
                "titles": "|".join(file_titles[:3]),
                "prop": "imageinfo",
                "iiprop": "url",
                "iiurlwidth": 330,
                "format": "json",
            }
            r = requests.get(search_url, params=params, headers={"User-Agent": USER_AGENT}, timeout=15)
            r.raise_for_status()
            data = r.json()

            for page_id, page in data.get("query", {}).get("pages", {}).items():
                if page_id == "-1":
                    continue
                imageinfo = page.get("imageinfo", [])
                if not imageinfo:
                    continue
                thumb_url = imageinfo[0].get("thumburl") or imageinfo[0].get("url")
                if thumb_url:
                    base = sanitize_filename(wiki_title)
                    save_path = os.path.join(IMG_DIR, base + "_commons")
                    result = _download_image(thumb_url, save_path)
                    if isinstance(result, str):
                        return result, "commons"
                    if result == 429:
                        time.sleep(2)
                        continue
            return None, "no usable image"
        except Exception as e:
            if attempt < 2:
                time.sleep(2)
                continue
            return None, f"{type(e).__name__}: {e}"
    return None, "max retries"


def _fetch_from_flickr(location_name):
    queries = [location_name, _ascii_fold(location_name)]
    words = location_name.replace("-", " ").split()
    if len(words) > 1:
        queries.append(words[-1])
        queries.append(" ".join(words[:2]))

    seen_urls = set()
    for attempt in range(5):
        query = queries[attempt] if attempt < len(queries) else queries[-1]
        try:
            url = "https://api.flickr.com/services/feeds/photos_public.gne"
            params = {
                "tags": query,
                "tagmode": "any",
                "format": "json",
                "nojsoncallback": 1,
                "lang": "en-us",
            }
            r = requests.get(url, params=params, headers={"User-Agent": USER_AGENT}, timeout=15)
            if r.status_code == 429:
                time.sleep(2)
                continue
            r.raise_for_status()
            data = r.json()
            items = data.get("items", [])
            base = sanitize_filename(location_name) + "_flickr"
            save_path = os.path.join(IMG_DIR, base)
            for item in items:
                img_url = item.get("media", {}).get("m")
                if not img_url or img_url in seen_urls:
                    continue
                seen_urls.add(img_url)
                for try_url in (img_url.replace("_m.jpg", ".jpg"), img_url):
                    result = _download_image(try_url, save_path)
                    if isinstance(result, str):
                        return result, "flickr"
                    if result != 429:
                        break
            if attempt >= min(len(queries), 3):
                return None, "no photos found"
        except Exception as e:
            if attempt < min(len(queries), 3):
                continue
            return None, f"{type(e).__name__}: {e}"
    return None, "max retries"


def _fetch_from_flickr_geo(lat, lon, location_name):
    for attempt in range(2):
        try:
            url = "https://api.flickr.com/services/feeds/photos_public.gne"
            params = {
                "lat": lat,
                "lon": lon,
                "radius": 1.0,
                "radius_units": "km",
                "format": "json",
                "nojsoncallback": 1,
                "lang": "en-us",
            }
            r = requests.get(url, params=params, headers={"User-Agent": USER_AGENT}, timeout=15)
            if r.status_code == 429:
                time.sleep(3)
                continue
            r.raise_for_status()
            data = r.json()
            items = data.get("items", [])
            if not items:
                return None, "no photos near location"
            base = sanitize_filename(location_name) + "_geo"
            save_path = os.path.join(IMG_DIR, base)
            for item in items[:6]:
                img_url = item.get("media", {}).get("m")
                if not img_url:
                    continue
                for try_url in (img_url.replace("_m.jpg", ".jpg"), img_url):
                    result = _download_image(try_url, save_path)
                    if isinstance(result, str):
                        return result, "flickr_geo"
                    if result != 429:
                        break
            return None, "no usable photos near location"
        except Exception as e:
            if attempt < 2:
                time.sleep(2)
                continue
            return None, f"{type(e).__name__}: {e}"
    return None, "max retries"


IMAGE_GAP_REPORT = []


def _report_missing_image(name, wiki_title):
    IMAGE_GAP_REPORT.append({
        "name": name,
        "wiki_title": wiki_title,
        "hint": f"https://en.wikipedia.org/wiki/{quote(wiki_title) if wiki_title else ''}" if wiki_title else None
    })
    return None


def _find_cached_image(name):
    needle = _ascii_fold(sanitize_filename(name)).lower()
    if not os.path.isdir(IMG_DIR):
        return None
    suffixes = ("_flickr", "_geo", "_commons")
    for f in os.listdir(IMG_DIR):
        f_noext, _ = os.path.splitext(f)
        for suffix in suffixes:
            if f_noext.endswith(suffix):
                f_noext = f_noext[:-len(suffix)]
                break
        f_folded = _ascii_fold(f_noext).lower()
        if f_folded == needle:
            return os.path.join(IMG_DIR, f)
    return None


def _resolve_manual_image(source, name):
    if source.startswith("http://") or source.startswith("https://"):
        save_path = os.path.join(IMG_DIR, sanitize_filename(name) + "_manual")
        result = _download_image(source, save_path)
        if isinstance(result, str):
            print(f"  [OK]   {name} (from manual URL)")
            return result
        print(f"  [..]   {name}: manual URL download failed")
        return None

    full_path = source if os.path.isabs(source) else os.path.join(SCRIPT_DIR, source)
    if os.path.exists(full_path):
        print(f"  [OK]   {name} (from manual file)")
        return full_path

    print(f"  [..]   {name}: manual source not found ({source})")
    return None


def fetch_image(wiki_title, location_name, lat, lon, image_source=None):
    name = location_name or wiki_title or ""

    if image_source:
        path = _resolve_manual_image(image_source, name)
        if path:
            return path

    existing = _find_cached_image(name)
    if existing:
        return existing

    path, reason = None, ""

    if wiki_title:
        path, reason = _fetch_from_wikipedia(wiki_title, name)
        if path:
            print(f"  [OK]   {name} (from Wikipedia)")
            return path
        print(f"  [..]   {name}: Wikipedia {reason}")

    if wiki_title:
        path, reason = _fetch_from_commons(wiki_title)
        if path:
            print(f"  [OK]   {name} (from Wikimedia Commons)")
            return path
        print(f"  [..]   {name}: Commons {reason}")

    path, reason = _fetch_from_flickr(name)
    if path:
        print(f"  [OK]   {name} (from Flickr)")
        return path
    print(f"  [..]   {name}: Flickr {reason}")

    if lat is not None and lon is not None:
        path, reason = _fetch_from_flickr_geo(lat, lon, name)
        if path:
            print(f"  [OK]   {name} (from Flickr geo)")
            return path
        print(f"  [..]   {name}: Flickr geo {reason}")

    _report_missing_image(name, wiki_title)
    print(f"  [GAP]  {name}: no image found, added to report")
    return None


def fetch_all_images(locations):
    futures = {}
    seen = set()
    with ThreadPoolExecutor(max_workers=4) as pool:
        for loc in locations:
            name = loc["name"]
            if name in seen:
                continue
            seen.add(name)
            wiki = loc.get("wiki")
            lat = loc.get("lat")
            lon = loc.get("lon")
            futures[pool.submit(fetch_image, wiki, name, lat, lon, loc.get("image"))] = name
        results = {}
        for fut in as_completed(futures):
            name = futures[fut]
            try:
                results[name] = fut.result()
            except Exception:
                results[name] = None
    return results


def print_image_gap_report():
    if not IMAGE_GAP_REPORT:
        return
    print()
    print("=" * 72)
    print("  IMAGE GAP REPORT — locations without any image source")
    print("=" * 72)
    for gap in IMAGE_GAP_REPORT:
        print(f"  - {gap['name']}")
        if gap.get("hint"):
            print(f"    Try: {gap['hint']}")
        print()
    print("  To fill gaps: download a suitable image and place it in")
    print(f"  the '{IMG_DIR}/' folder with the filename matching the")
    print("  location name (e.g. 'Thorufoss.jpg'). Then rerun the script.")
    print("=" * 72)
    print()

# ---------------------------------------------------------------------------
# Map helpers
# ---------------------------------------------------------------------------

def make_map(day_idx, locs):
    path = os.path.join(MAP_DIR, "map_{}.png".format(day_idx))
    m = StaticMap(800, 800, url_template="https://a.tile.openstreetmap.org/{z}/{x}/{y}.png")
    coords = [(loc["lon"], loc["lat"]) for loc in locs if loc.get("lat") and loc.get("lon")]
    if len(coords) > 1:
        m.add_line(Line(coords, "blue", 3))
    for lon, lat in coords:
        m.add_marker(CircleMarker((lon, lat), "red", 12))
        m.add_marker(CircleMarker((lon, lat), "white", 6))
    img = m.render()
    img.save(path)
    return path


def gmaps_link(locs):
    pts = "/".join("{},{}".format(loc["lat"], loc["lon"]) for loc in locs
                   if loc.get("lat") and loc.get("lon"))
    return "https://www.google.com/maps/dir/" + pts


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Day planning
# ---------------------------------------------------------------------------

WEEKDAYS = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]


def haversine(lat1, lon1, lat2, lon2):
    R = 6371
    dlat = math.radians(lat2 - lat1)
    dlon = math.radians(lon2 - lon1)
    a = math.sin(dlat / 2) ** 2 + math.cos(math.radians(lat1)) * math.cos(math.radians(lat2)) * math.sin(dlon / 2) ** 2
    return R * 2 * math.atan2(math.sqrt(a), math.sqrt(1 - a))


def _route_distance(locs):
    dist = 0
    for i in range(1, len(locs)):
        if locs[i].get("lat") and locs[i - 1].get("lat"):
            dist += haversine(locs[i - 1]["lat"], locs[i - 1]["lon"],
                              locs[i]["lat"], locs[i]["lon"])
    return dist


def _format_date(date_str):
    dt = datetime.strptime(date_str, "%Y-%m-%d")
    return f"{WEEKDAYS[dt.weekday()]} {dt.day} June"


def _fmt_visit(hours):
    if hours < 0.75:
        return f"{int(hours * 60)}min"
    h = int(hours)
    m = int(round((hours - h) * 60))
    return f"{h}h{m}min" if m else f"{h}h"


def plan_days(data):
    trip = data["trip"]
    locations = [loc for loc in data["locations"] if loc.get("is_included")]

    airport = trip["airport"]
    max_hours = trip.get("max_day_hours", 10)
    avg_speed = trip.get("avg_speed_kmh", 70)

    bookings = sorted(
        [loc for loc in locations if loc.get("is_booking") and loc.get("start_booking")],
        key=lambda x: x["start_booking"]
    )
    attractions = [loc for loc in locations
                   if not loc.get("is_booking") and loc.get("type") != "transport"]

    # Build home-base periods from consecutive bookings
    periods = []
    for i, b in enumerate(bookings):
        end_date = bookings[i + 1]["start_booking"] if i + 1 < len(bookings) else trip["end_date"]
        periods.append({
            "base": b,
            "date_start": b["start_booking"],
            "date_end": end_date,
        })

    # Assign attractions based on geographic domain + tiebreak by Haversine
    # Trip route: Reykjavik area -> Golden Circle -> Reykjanes ->
    #             Snaefellsnes (north) -> South Coast -> Glaciers
    for sight in attractions:
        lat, lon = sight["lat"], sight["lon"]
        # Manual domain override in JSON, else algorithmic assignment
        domain = sight.get("domain")
        if domain is None:
            if lat >= 64.5 and len(periods) > 1:
                domain = 1
            elif lon > -20.0 and lat < 64.5 and len(periods) > 2:
                domain = 2
            else:
                domain = 0

        best_period = periods[min(domain, len(periods) - 1)]
        best_period.setdefault("sights", []).append(sight)

    for p in periods:
        if "sights" not in p:
            p["sights"] = []
        p["sights"].sort(key=lambda s: haversine(s["lat"], s["lon"],
                                                  p["base"]["lat"], p["base"]["lon"]))

    # -----------------------------------------------------------------------
    # Build daily groups per period using time-based clustering
    # -----------------------------------------------------------------------
    days_out = []
    day_idx = 1

    # Day 1: airport -> first accommodation (no sightseeing)
    day_title = f"Day {day_idx} - {_format_date(trip['start_date'])} - Arrival"
    days_out.append({
        "title": day_title,
        "locations": [airport, bookings[0]] if bookings else [airport],
        "date": trip["start_date"],
    })
    day_idx += 1

    for period_idx, p in enumerate(periods):
        p_start = datetime.strptime(p["date_start"], "%Y-%m-%d")
        p_end = datetime.strptime(p["date_end"], "%Y-%m-%d")
        num_days = (p_end - p_start).days

        sights = list(p["sights"])
        if not sights:
            continue

        # Skip j=0 for the first period (arrival day already used)
        start_offset = 1 if p == periods[0] else 0
        available_slots = num_days - start_offset
        if available_slots <= 0:
            continue

        # Sort sights by longitude (west to east) and split into contiguous chunks
        sorted_sights = sorted(sights, key=lambda s: -s["lon"])
        n_slots = min(len(sorted_sights), available_slots)
        groups = []
        for i in range(n_slots):
            start = i * len(sorted_sights) // n_slots
            end = (i + 1) * len(sorted_sights) // n_slots
            groups.append(sorted_sights[start:end])

        # Apply manual day overrides
        for s in sights:
            d = s.get("day")
            if d is not None and 0 <= d < n_slots:
                for g in groups:
                    if s in g and g is not groups[d]:
                        g.remove(s)
                        groups[d].append(s)
                        break

        # Assign groups to dates, with start/end at base
        # TSP ordering within each group starts from the day's start base
        for j, group in enumerate(groups):
            date = p_start + timedelta(days=start_offset + j)
            date_str = date.strftime("%Y-%m-%d")

            if period_idx > 0 and j == 0:
                start_base = periods[period_idx - 1]["base"]
                end_base = p["base"]
            else:
                start_base = p["base"]
                end_base = p["base"]

            # Nearest-neighbour TSP starting from start_base
            ordered = [min(group, key=lambda s: haversine(start_base["lat"], start_base["lon"], s["lat"], s["lon"]))]
            remaining = [s for s in group if s is not ordered[0]]
            while remaining:
                last = ordered[-1]
                best = min(remaining, key=lambda s: haversine(last["lat"], last["lon"], s["lat"], s["lon"]))
                ordered.append(best)
                remaining.remove(best)

            route = [start_base] + ordered + [end_base]
            route_km = int(_route_distance(route))
            total_time = sum(loc.get("est_visit_time", 0.5) for loc in group) + route_km / avg_speed
            day_title = f"Day {day_idx} - {_format_date(date_str)} ({route_km} km, {total_time:.1f}h)"

            days_out.append({
                "title": day_title,
                "locations": route,
                "date": date_str,
            })
            day_idx += 1

    # -----------------------------------------------------------------------
    # Final day: return to airport (separate day for the drive back)
    # -----------------------------------------------------------------------
    last_date = days_out[-1].get("date", "")
    if last_date != trip["end_date"]:
        last_accommodation = bookings[-1] if bookings else airport
        route = [last_accommodation, airport]
        route_km = int(_route_distance(route))
        day_title = f"Day {day_idx} - {_format_date(trip['end_date'])} ({route_km} km) - Return to Airport"
        days_out.append({
            "title": day_title,
            "locations": route,
            "date": trip["end_date"],
        })
    else:
        last = days_out[-1]
        if airport not in last["locations"]:
            last["locations"].append(airport)
        last["title"] = last["title"].replace(" (", " → Airport (")

    return days_out

# ---------------------------------------------------------------------------
# PPTX building
# ---------------------------------------------------------------------------

def build_presentation(days):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = find_blank_layout(prs)

    for i, day in enumerate(days):
        title = day["title"]
        locs = day["locations"]
        s = prs.slides.add_slide(blank_layout)

        tb = s.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x3B, 0x73)

        try:
            mp = make_map(i, locs)
            s.shapes.add_picture(mp, Inches(7.4), Inches(1.0), height=Inches(5.6))
        except Exception as e:
            print("map error:", e)

        lb = s.shapes.add_textbox(Inches(7.4), Inches(6.7), Inches(5.8), Inches(0.5))
        run = lb.text_frame.paragraphs[0].add_run()
        run.text = "Open route in Google Maps"
        run.font.size = Pt(12)
        run.hyperlink.address = gmaps_link(locs)

        images = fetch_all_images(locs)

        top = 1.0
        for loc in locs:
            name = loc["name"]
            desc = loc.get("desc", "")
            img = images.get(name)

            desc_lines = max(1, math.ceil(len(desc) / 105))
            row_height = 0.22 + desc_lines * 0.18 + 0.10

            if img:
                try:
                    s.shapes.add_picture(img, Inches(0.3), Inches(top), height=Inches(0.55))
                except Exception:
                    pass
            elif not img:
                try:
                    icon = s.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(top), Inches(0.55), Inches(0.55)
                    )
                    icon.fill.solid()
                    icon.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    icon.line.fill.background()
                except Exception:
                    pass

            box = s.shapes.add_textbox(Inches(1.15), Inches(top), Inches(6.0), Inches(row_height))
            tf = box.text_frame
            tf.word_wrap = True
            r1 = tf.paragraphs[0].add_run()
            r1.text = name
            r1.font.bold = True
            r1.font.size = Pt(13)
            r1.font.color.rgb = RGBColor(0x1F, 0x3B, 0x73)
            v = loc.get("est_visit_time")
            if v:
                r1b = tf.paragraphs[0].add_run()
                r1b.text = "  \u2014  " + _fmt_visit(v)
                r1b.font.size = Pt(11)
                r1b.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            para = tf.add_paragraph()
            r2 = para.add_run()
            r2.text = desc
            r2.font.size = Pt(10.5)
            top += row_height + 0.06

    prs.save(OUTPUT_FILE)
    print("Saved", OUTPUT_FILE)


def main():
    IMAGE_GAP_REPORT.clear()

    with open(DATA_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)

    days = plan_days(data)

    print(f"\nPlanned {len(days)} days from {len(data['locations'])} locations\n")

    build_presentation(days)

    print_image_gap_report()


if __name__ == "__main__":
    main()
